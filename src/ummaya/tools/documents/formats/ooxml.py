# SPDX-License-Identifier: Apache-2.0
"""OOXML adapter and engine boundaries for DOCX, XLSX, and PPTX."""

from __future__ import annotations

import html
import re
from copy import copy
from datetime import date, datetime
from decimal import Decimal
from pathlib import Path
from typing import TYPE_CHECKING, BinaryIO, Protocol, cast

import docx
import openpyxl  # type: ignore[import-untyped]
import pptx
from docx.document import Document as DocxDocument
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from docx.table import Table as DocxTable
from docx.table import _Cell as DocxCell
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.text.run import Run as DocxRun
from openpyxl.cell.cell import Cell  # type: ignore[import-untyped]
from openpyxl.styles import Alignment, Font, PatternFill  # type: ignore[import-untyped]
from openpyxl.worksheet.worksheet import Worksheet  # type: ignore[import-untyped]
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PptxPresentation
from pptx.slide import Slide as PptxSlide
from pptx.table import Table as PptxTable

from ummaya.tools.documents.engines import DocumentInspectionEngine, DocumentMutationEngine
from ummaya.tools.documents.models import (
    DocumentExtraction,
    DocumentFormat,
    DocumentPatch,
    DocumentPatchOperation,
    ImageReference,
    KnownDocumentFormat,
    MetadataValue,
    OperationType,
    ParagraphBlock,
    ScalarValue,
    StyleDescriptor,
    TableBlock,
    TableCell,
)

if TYPE_CHECKING:
    from ummaya.tools.documents.tool_defs import DocumentFieldPatch


class _OfficeSaveable(Protocol):
    """Document object that can save itself to a binary file-like stream."""

    def save(self, target: BinaryIO) -> None:
        """Persist the office document into the provided binary stream."""


OOXML_CANDIDATE_ENGINES: dict[DocumentFormat, tuple[str, ...]] = {
    DocumentFormat.docx: ("python-docx", "direct-wordprocessingml-oracle"),
    DocumentFormat.xlsx: ("openpyxl", "direct-spreadsheetml-oracle"),
    DocumentFormat.pptx: ("python-pptx", "direct-presentationml-oracle"),
}

_DOCX_PARAGRAPH_RE = re.compile(r"(?:^|/)paragraphs?/(?P<paragraph>\d+)(?:/runs/(?P<run>\d+))?$")
_DOCX_TABLE_CELL_RE = re.compile(
    r"(?:^|/)tables?/(?P<table>\d+)/rows?/(?P<row>\d+)/cells?/(?P<cell>\d+)$|"
    r"(?:^|/)table/(?P<table2>\d+)/r(?P<row2>\d+)c(?P<cell2>\d+)$"
)
_XLSX_CELL_RE = re.compile(r"^/sheets/(?P<sheet>[^/]+)/cells/(?P<cell>[A-Za-z]{1,3}\d+)$")
_PPTX_SHAPE_TEXT_RE = re.compile(r"^/slides/(?P<slide>\d+)/shapes/(?P<shape>\d+)/text$")
_PPTX_TABLE_CELL_RE = re.compile(
    r"^/slides/(?P<slide>\d+)/tables/(?P<table>\d+)/rows/(?P<row>\d+)/cells/(?P<cell>\d+)$"
)


def validate_ooxml_engine(engine: DocumentInspectionEngine) -> DocumentInspectionEngine:
    """Validate that an injected engine is scoped to an OOXML format."""
    if engine.document_format not in OOXML_CANDIDATE_ENGINES:
        raise ValueError("OOXML adapter requires a docx, xlsx, or pptx engine")
    return engine


def validate_ooxml_mutation_engine(engine: DocumentInspectionEngine) -> DocumentMutationEngine:
    """Validate that an injected OOXML engine can safely mutate derivatives."""
    validate_ooxml_engine(engine)
    if not isinstance(engine, DocumentMutationEngine):
        raise ValueError("OOXML adapter requires a mutation-capable engine")
    return engine


class _OoxmlAdapterBase:
    """Shared adapter behavior for one OOXML file family."""

    adapter_id: str
    known_formats: tuple[KnownDocumentFormat, ...]
    promoted_formats: tuple[DocumentFormat, ...]

    def __init__(self, inspection_engine: DocumentInspectionEngine | None = None) -> None:
        self._inspection_engine = (
            validate_ooxml_engine(inspection_engine) if inspection_engine is not None else None
        )

    @property
    def engine_id(self) -> str:
        """Return the wrapped engine id for diagnostics."""
        if self._inspection_engine is None:
            return self.adapter_id
        return self._inspection_engine.engine_id

    def inspect(self, path: Path, *, artifact_id: str) -> DocumentExtraction:
        """Delegate inspection when promoted; otherwise return a typed empty scope."""
        if self._inspection_engine is None:
            return DocumentExtraction(
                artifact_id=artifact_id,
                metadata={"adapter_id": self.adapter_id},
                warnings=[f"{self.adapter_id} is registered as known-only."],
            )
        return self._inspection_engine.inspect(path, artifact_id=artifact_id)

    def normalize_fill_patches(
        self,
        patches: tuple[DocumentFieldPatch, ...],
        *,
        extraction: DocumentExtraction | None,
    ) -> tuple[DocumentFieldPatch, ...]:
        """Return fill patches unchanged for OOXML adapters."""
        _ = extraction
        return patches


class DocxDocumentAdapter(_OoxmlAdapterBase):
    """DOCX adapter boundary backed by python-docx."""

    adapter_id: str = "python-docx-adapter"
    known_formats: tuple[KnownDocumentFormat, ...] = (KnownDocumentFormat.docx,)
    promoted_formats: tuple[DocumentFormat, ...] = (DocumentFormat.docx,)

    def __init__(self, inspection_engine: DocumentInspectionEngine | None = None) -> None:
        super().__init__(inspection_engine or PythonDocxDocumentEngine())


class XlsxDocumentAdapter(_OoxmlAdapterBase):
    """XLSX adapter boundary backed by openpyxl when promoted."""

    adapter_id: str = "openpyxl-adapter"
    known_formats: tuple[KnownDocumentFormat, ...] = (KnownDocumentFormat.xlsx,)
    promoted_formats: tuple[DocumentFormat, ...] = ()

    def __init__(self, inspection_engine: DocumentInspectionEngine | None = None) -> None:
        self.promoted_formats = (DocumentFormat.xlsx,) if inspection_engine is not None else ()
        super().__init__(inspection_engine)


class PptxDocumentAdapter(_OoxmlAdapterBase):
    """PPTX adapter boundary backed by python-pptx when promoted."""

    adapter_id: str = "python-pptx-adapter"
    known_formats: tuple[KnownDocumentFormat, ...] = (KnownDocumentFormat.pptx,)
    promoted_formats: tuple[DocumentFormat, ...] = ()

    def __init__(self, inspection_engine: DocumentInspectionEngine | None = None) -> None:
        self.promoted_formats = (DocumentFormat.pptx,) if inspection_engine is not None else ()
        super().__init__(inspection_engine)


class PythonDocxDocumentEngine:
    """DOCX read/write engine backed by the promoted python-docx dependency."""

    document_format = DocumentFormat.docx
    engine_id = "python-docx"
    render_artifact_extension = "svg"
    render_mime_type = "image/svg+xml"

    def inspect(self, path: Path, *, artifact_id: str) -> DocumentExtraction:
        """Extract normalized paragraphs, tables, and core metadata from DOCX."""
        document = docx.Document(str(path))
        paragraphs: list[ParagraphBlock] = []
        tables: list[TableBlock] = []

        paragraph_index = 1
        table_index = 1
        for block in document.iter_inner_content():
            if isinstance(block, DocxParagraph):
                if block.text:
                    paragraphs.append(
                        _paragraph_block(
                            block,
                            engine_id=self.engine_id,
                            path=path,
                            index=paragraph_index,
                        )
                    )
                    paragraph_index += 1
            elif isinstance(block, DocxTable):
                tables.append(
                    _table_block(
                        block,
                        engine_id=self.engine_id,
                        path=path,
                        index=table_index,
                    )
                )
                table_index += 1

        return DocumentExtraction(
            artifact_id=artifact_id,
            paragraphs=paragraphs,
            tables=tables,
            metadata=_docx_core_metadata(document),
            warnings=[
                "python-docx scope excludes nested tables and tracked revision "
                "content from the top-level document lists."
            ],
        )

    def apply_patch(self, path: Path, patch: DocumentPatch) -> bytes:
        """Apply bounded paragraph, run, table-cell, style, and metadata edits."""
        document = docx.Document(str(path))
        for operation in patch.operations:
            _apply_docx_operation(document, operation)
        output = _save_to_bytes(document)
        return output

    def render(self, path: Path, *, artifact_id: str, output_dir: Path) -> tuple[bytes, ...]:
        """Render a lightweight SVG evidence page for DOCX review."""
        _ = output_dir
        extraction = self.inspect(path, artifact_id=artifact_id)
        lines = [block.text for block in extraction.paragraphs]
        lines.extend(cell.text for table in extraction.tables for cell in table.cells if cell.text)
        return (_svg_page(lines or [Path(path).name], title=f"DOCX {artifact_id}"),)


PythonDocxInspectionEngine = PythonDocxDocumentEngine


class OpenPyxlDocumentEngine:
    """XLSX read/write engine backed by openpyxl."""

    document_format = DocumentFormat.xlsx
    engine_id = "openpyxl"
    render_artifact_extension = "svg"
    render_mime_type = "image/svg+xml"

    def inspect(self, path: Path, *, artifact_id: str) -> DocumentExtraction:
        """Extract normalized sheet cells, metadata, and style anchors."""
        workbook = openpyxl.load_workbook(path, data_only=False)
        tables: list[TableBlock] = []
        style_map: list[StyleDescriptor] = []
        metadata: dict[str, MetadataValue] = {
            "engine_id": self.engine_id,
            "format": "xlsx",
            "sheet_count": len(workbook.worksheets),
        }
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            cells: list[TableCell] = []
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    source_path = f"/sheets/{worksheet.title}/cells/{cell.coordinate}"
                    cells.append(
                        TableCell(
                            row_index=cell.row - 1,
                            column_index=cell.column - 1,
                            text=str(cell.value),
                            source_path=source_path,
                            field_path=source_path,
                        )
                    )
                    if _cell_has_non_default_style(cell):
                        style_map.append(_xlsx_style_descriptor(cell, source_path))
            tables.append(
                TableBlock(
                    block_id=f"xlsx-sheet-{sheet_index:03d}",
                    source_path=f"/sheets/{worksheet.title}",
                    cells=cells,
                )
            )
            if worksheet.print_area:
                metadata[f"sheet_{sheet_index}_print_area"] = str(worksheet.print_area)

        return DocumentExtraction(
            artifact_id=artifact_id,
            tables=tables,
            metadata=metadata,
            style_map=style_map,
            warnings=[
                "openpyxl preserves formula strings but UMMAYA does not claim formula "
                "evaluation or cached-value recalculation."
            ],
        )

    def apply_patch(self, path: Path, patch: DocumentPatch) -> bytes:
        """Apply bounded cell, cell-style, and workbook metadata edits."""
        workbook = openpyxl.load_workbook(path, data_only=False)
        for operation in patch.operations:
            _apply_xlsx_operation(workbook, operation)
        return _save_workbook_to_bytes(workbook)

    def render(self, path: Path, *, artifact_id: str, output_dir: Path) -> tuple[bytes, ...]:
        """Render one SVG evidence page per worksheet."""
        _ = output_dir
        extraction = self.inspect(path, artifact_id=artifact_id)
        pages: list[bytes] = []
        for table in extraction.tables:
            lines = [cell.text for cell in table.cells[:36]]
            pages.append(_svg_page(lines or [table.source_path], title=table.block_id))
        return tuple(pages) or (_svg_page([Path(path).name], title=f"XLSX {artifact_id}"),)


class PythonPptxDocumentEngine:
    """PPTX read/write engine backed by python-pptx."""

    document_format = DocumentFormat.pptx
    engine_id = "python-pptx"
    render_artifact_extension = "svg"
    render_mime_type = "image/svg+xml"

    def inspect(self, path: Path, *, artifact_id: str) -> DocumentExtraction:
        """Extract normalized slide text, tables, images, and core metadata."""
        presentation = pptx.Presentation(str(path))
        paragraphs: list[ParagraphBlock] = []
        tables: list[TableBlock] = []
        images: list[ImageReference] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            table_index = 1
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_text_frame", False) and shape.text:
                    paragraphs.append(
                        ParagraphBlock(
                            block_id=f"pptx-slide-{slide_index:03d}-shape-{shape_index:03d}",
                            text=shape.text,
                            source_path=f"/slides/{slide_index}/shapes/{shape_index}/text",
                        )
                    )
                if getattr(shape, "has_table", False):
                    tables.append(
                        _pptx_table_block(
                            cast(PptxTable, shape.table),
                            slide_index=slide_index,
                            table_index=table_index,
                        )
                    )
                    table_index += 1
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(
                        ImageReference(
                            image_id=f"pptx-slide-{slide_index:03d}-image-{shape_index:03d}",
                            source_path=f"/slides/{slide_index}/images/{shape_index}",
                            content_type=getattr(shape.image, "content_type", "image/unknown"),
                            alt_text=getattr(shape, "name", None),
                        )
                    )

        return DocumentExtraction(
            artifact_id=artifact_id,
            paragraphs=paragraphs,
            tables=tables,
            images=images,
            metadata=_pptx_core_metadata(presentation),
            warnings=[
                "python-pptx scope blocks animations, masters, and complex media rewrites "
                "until separate promotion gates pass."
            ],
        )

    def apply_patch(self, path: Path, patch: DocumentPatch) -> bytes:
        """Apply bounded slide text, table-cell, and core metadata edits."""
        presentation = pptx.Presentation(str(path))
        for operation in patch.operations:
            _apply_pptx_operation(presentation, operation)
        return _save_to_bytes(presentation)

    def render(self, path: Path, *, artifact_id: str, output_dir: Path) -> tuple[bytes, ...]:
        """Render one SVG evidence page per slide."""
        _ = output_dir
        extraction = self.inspect(path, artifact_id=artifact_id)
        by_slide: dict[int, list[str]] = {}
        for block in extraction.paragraphs:
            slide_index = _slide_index_from_path(block.source_path)
            by_slide.setdefault(slide_index, []).append(block.text)
        for table in extraction.tables:
            slide_index = _slide_index_from_path(table.source_path)
            by_slide.setdefault(slide_index, []).extend(cell.text for cell in table.cells)
        return tuple(
            _svg_page(lines or [f"Slide {slide_index}"], title=f"PPTX slide {slide_index}")
            for slide_index, lines in sorted(by_slide.items())
        ) or (_svg_page([Path(path).name], title=f"PPTX {artifact_id}"),)


def _paragraph_block(
    paragraph: DocxParagraph,
    *,
    engine_id: str,
    path: Path,
    index: int,
) -> ParagraphBlock:
    return ParagraphBlock(
        block_id=f"docx-paragraph-{index:03d}",
        text=paragraph.text,
        source_path=f"engine://{engine_id}/{path.name}/paragraph/{index}",
        style_id=_paragraph_style_id(paragraph),
    )


def _table_block(
    table: DocxTable,
    *,
    engine_id: str,
    path: Path,
    index: int,
) -> TableBlock:
    cells: list[TableCell] = []
    for row_index, row in enumerate(table.rows):
        row_cells = tuple(row.cells)
        for column_index, cell in enumerate(row_cells):
            source_path = (
                f"engine://{engine_id}/{path.name}/table/{index}/"
                f"r{row_index + 1}c{column_index + 1}"
            )
            cells.append(
                TableCell(
                    row_index=row_index,
                    column_index=column_index,
                    text=cell.text,
                    source_path=source_path,
                    field_path=(
                        source_path
                        if _docx_adjacent_label_blank_value_cell(row_cells, column_index)
                        else None
                    ),
                )
            )
    return TableBlock(
        block_id=f"docx-table-{index:03d}",
        source_path=f"engine://{engine_id}/{path.name}/table/{index}",
        cells=cells,
    )


def _docx_adjacent_label_blank_value_cell(
    row_cells: tuple[DocxCell, ...],
    column_index: int,
) -> bool:
    if column_index <= 0:
        return False
    if row_cells[column_index].text.strip():
        return False
    return _docx_meaningful_form_label(row_cells[column_index - 1].text)


def _docx_meaningful_form_label(text: str) -> bool:
    normalized = re.sub(r"\s+", "", text)
    if len(normalized) < 2:
        return False
    return re.search(r"[0-9A-Za-z가-힣]", normalized) is not None


def _paragraph_style_id(paragraph: DocxParagraph) -> str | None:
    style: object | None = paragraph.style
    style_id = getattr(style, "style_id", None)
    if isinstance(style_id, str) and style_id:
        return style_id
    style_name = getattr(style, "name", None)
    if isinstance(style_name, str) and style_name:
        return style_name
    return None


def _docx_core_metadata(document: DocxDocument) -> dict[str, MetadataValue]:
    core_properties = document.core_properties
    metadata: dict[str, MetadataValue] = {
        "engine_id": "python-docx",
        "format": "docx",
    }
    for property_name in (
        "author",
        "category",
        "comments",
        "content_status",
        "created",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "last_printed",
        "modified",
        "revision",
        "subject",
        "title",
        "version",
    ):
        value = getattr(core_properties, property_name)
        if _metadata_value_is_present(value):
            metadata[f"core_{property_name}"] = value
    return metadata


def _metadata_value_is_present(value: MetadataValue) -> bool:
    if isinstance(value, str):
        return bool(value)
    if isinstance(value, datetime):
        return True
    return value is not None


def _apply_docx_operation(document: DocxDocument, operation: DocumentPatchOperation) -> None:
    if operation.operation_type in {
        OperationType.replace_text,
        OperationType.set_field_value,
        OperationType.insert_paragraph,
    }:
        _apply_docx_text_operation(document, operation)
        return
    if operation.operation_type is OperationType.set_table_cell:
        table, row_index, cell_index = _docx_table_cell(document, operation.target_path)
        _ = table
        cell = document.tables[_docx_table_ordinal(operation.target_path) - 1].cell(
            row_index,
            cell_index,
        )
        _set_docx_paragraph_text(cell.paragraphs[0], _string_value(operation.value))
        return
    if operation.operation_type is OperationType.set_document_metadata:
        _set_docx_metadata(document, operation)
        return
    if operation.operation_type is OperationType.set_run_style:
        paragraph_index, run_index = _docx_paragraph_and_run_indexes(operation.target_path)
        _apply_docx_run_style(document.paragraphs[paragraph_index], run_index, operation.style)
        return
    if operation.operation_type is OperationType.set_paragraph_style:
        paragraph_index, _ = _docx_paragraph_and_run_indexes(operation.target_path)
        _apply_docx_paragraph_style(document.paragraphs[paragraph_index], operation.style)
        return
    if operation.operation_type is OperationType.set_cell_style:
        table, row_index, cell_index = _docx_table_cell(document, operation.target_path)
        cell = table.cell(row_index, cell_index)
        _apply_docx_cell_style(cell, operation.style)
        return
    raise ValueError(f"Unsupported DOCX operation: {operation.operation_type.value}")


def _apply_docx_text_operation(
    document: DocxDocument,
    operation: DocumentPatchOperation,
) -> None:
    paragraph_index, run_index = _docx_paragraph_and_run_indexes(operation.target_path)
    paragraph = document.paragraphs[paragraph_index]
    if operation.operation_type is OperationType.insert_paragraph:
        document.add_paragraph(_string_value(operation.value))
        return
    if run_index is None:
        _set_docx_paragraph_text(paragraph, _string_value(operation.value))
        return
    while len(paragraph.runs) <= run_index:
        paragraph.add_run("")
    paragraph.runs[run_index].text = _string_value(operation.value)


def _set_docx_paragraph_text(paragraph: DocxParagraph, value: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = value
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run(value)


def _docx_paragraph_and_run_indexes(target_path: str) -> tuple[int, int | None]:
    match = _DOCX_PARAGRAPH_RE.search(target_path)
    if match is None:
        raise ValueError(f"Unsupported DOCX paragraph target: {target_path}")
    paragraph_index = int(match.group("paragraph")) - 1
    run_value = match.group("run")
    return paragraph_index, int(run_value) - 1 if run_value is not None else None


def _docx_table_ordinal(target_path: str) -> int:
    match = _DOCX_TABLE_CELL_RE.search(target_path)
    if match is None:
        raise ValueError(f"Unsupported DOCX table target: {target_path}")
    return int(match.group("table") or match.group("table2"))


def _docx_table_cell(
    document: DocxDocument,
    target_path: str,
) -> tuple[DocxTable, int, int]:
    match = _DOCX_TABLE_CELL_RE.search(target_path)
    if match is None:
        raise ValueError(f"Unsupported DOCX table target: {target_path}")
    table_index = int(match.group("table") or match.group("table2")) - 1
    row_index = int(match.group("row") or match.group("row2")) - 1
    cell_index = int(match.group("cell") or match.group("cell2")) - 1
    return document.tables[table_index], row_index, cell_index


def _set_docx_metadata(document: DocxDocument, operation: DocumentPatchOperation) -> None:
    property_name = operation.target_path.rsplit("/", 1)[-1]
    if property_name not in {
        "author",
        "category",
        "comments",
        "content_status",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "revision",
        "subject",
        "title",
        "version",
    }:
        raise ValueError(f"Unsupported DOCX core metadata target: {operation.target_path}")
    setattr(document.core_properties, property_name, _string_value(operation.value))


def _apply_docx_run_style(
    paragraph: DocxParagraph,
    run_index: int | None,
    style: StyleDescriptor | None,
) -> None:
    if run_index is None:
        raise ValueError("DOCX run style target must include /runs/{index}")
    if style is None:
        raise ValueError("DOCX run style operation requires style")
    while len(paragraph.runs) <= run_index:
        paragraph.add_run("")
    run = paragraph.runs[run_index]
    if style.bold is not None:
        run.bold = style.bold
    if style.italic is not None:
        run.italic = style.italic
    if style.underline is not None:
        run.underline = style.underline
    if style.font_family is not None:
        _set_docx_run_font_family(run, style.font_family)
    if style.font_size_pt is not None:
        run.font.size = Pt(float(style.font_size_pt))
    if style.font_color_rgb is not None:
        run.font.color.rgb = RGBColor.from_string(style.font_color_rgb)


def _apply_docx_paragraph_style(
    paragraph: DocxParagraph,
    style: StyleDescriptor | None,
) -> None:
    if style is None:
        raise ValueError("DOCX paragraph style operation requires style")
    _apply_docx_paragraph_alignment(paragraph, style)
    if _docx_style_has_direct_run_properties(style):
        _apply_docx_direct_style_to_paragraph_runs(paragraph, style)
        return
    try:
        paragraph.style = style.style_id
    except KeyError:
        return


def _apply_docx_cell_style(cell: DocxCell, style: StyleDescriptor | None) -> None:
    if style is None:
        raise ValueError("DOCX cell style operation requires style")
    if style.fill_color_rgb is not None:
        _apply_docx_cell_fill(cell, style.fill_color_rgb)
    for paragraph in cell.paragraphs:
        _apply_docx_paragraph_alignment(paragraph, style)
        if _docx_style_has_direct_run_properties(style):
            _apply_docx_direct_style_to_paragraph_runs(paragraph, style)


def _apply_docx_cell_fill(cell: DocxCell, fill_color_rgb: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()  # noqa: SLF001 - python-docx exposes cell shading only via OOXML.
    shading = tc_pr.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        tc_pr.append(shading)
    shading.set(qn("w:fill"), fill_color_rgb.upper())


def _apply_docx_paragraph_alignment(
    paragraph: DocxParagraph,
    style: StyleDescriptor,
) -> None:
    if style.alignment is None:
        return
    alignment = {
        "left": WD_ALIGN_PARAGRAPH.LEFT,
        "center": WD_ALIGN_PARAGRAPH.CENTER,
        "right": WD_ALIGN_PARAGRAPH.RIGHT,
        "justify": WD_ALIGN_PARAGRAPH.JUSTIFY,
        "distributed": WD_ALIGN_PARAGRAPH.DISTRIBUTE,
    }[style.alignment]
    paragraph.alignment = alignment


def _docx_style_has_direct_run_properties(style: StyleDescriptor) -> bool:
    return any(
        value is not None
        for value in (
            style.bold,
            style.italic,
            style.underline,
            style.font_family,
            style.font_size_pt,
            style.font_color_rgb,
        )
    )


def _apply_docx_direct_style_to_paragraph_runs(
    paragraph: DocxParagraph,
    style: StyleDescriptor,
) -> None:
    runs = paragraph.runs or [paragraph.add_run("")]
    for run in runs:
        _apply_docx_direct_run_style(run, style)


def _apply_docx_direct_run_style(run: DocxRun, style: StyleDescriptor) -> None:
    if style.bold is not None:
        run.bold = style.bold
    if style.italic is not None:
        run.italic = style.italic
    if style.underline is not None:
        run.underline = style.underline
    if style.font_family is not None:
        _set_docx_run_font_family(run, style.font_family)
    if style.font_size_pt is not None:
        run.font.size = Pt(float(style.font_size_pt))
    if style.font_color_rgb is not None:
        run.font.color.rgb = RGBColor.from_string(style.font_color_rgb)


def _set_docx_run_font_family(run: DocxRun, font_family: str) -> None:
    run.font.name = font_family
    r_pr = run._element.get_or_add_rPr()  # noqa: SLF001 - CJK fonts require raw run OOXML.
    r_fonts = r_pr.find(qn("w:rFonts"))
    if r_fonts is None:
        r_fonts = OxmlElement("w:rFonts")
        r_pr.insert(0, r_fonts)
    for attribute_name in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
        r_fonts.set(qn(attribute_name), font_family)


def _apply_xlsx_operation(workbook: openpyxl.Workbook, operation: DocumentPatchOperation) -> None:
    if operation.operation_type is OperationType.set_table_cell:
        worksheet, cell_ref = _xlsx_cell_target(workbook, operation.target_path)
        _ensure_xlsx_cell_editable(worksheet, cell_ref)
        worksheet[cell_ref].value = _office_scalar(operation.value)
        return
    if operation.operation_type is OperationType.set_cell_style:
        worksheet, cell_ref = _xlsx_cell_target(workbook, operation.target_path)
        _ensure_xlsx_cell_editable(worksheet, cell_ref)
        _apply_xlsx_cell_style(worksheet[cell_ref], operation.style)
        return
    if operation.operation_type is OperationType.set_document_metadata:
        _set_xlsx_metadata(workbook, operation)
        return
    raise ValueError(f"Unsupported XLSX operation: {operation.operation_type.value}")


def _xlsx_cell_target(
    workbook: openpyxl.Workbook,
    target_path: str,
) -> tuple[Worksheet, str]:
    match = _XLSX_CELL_RE.match(target_path)
    if match is None:
        raise ValueError(f"Unsupported XLSX cell target: {target_path}")
    sheet_name = match.group("sheet")
    if sheet_name not in workbook.sheetnames:
        raise ValueError(f"XLSX sheet does not exist: {sheet_name}")
    return workbook[sheet_name], match.group("cell").upper()


def _ensure_xlsx_cell_editable(worksheet: Worksheet, cell_ref: str) -> None:
    merged_ranges = worksheet.merged_cells.ranges
    for merged_range in merged_ranges:
        if cell_ref in merged_range and cell_ref != str(merged_range).split(":", 1)[0]:
            raise ValueError(f"Cannot edit non-anchor merged cell: {cell_ref}")


def _apply_xlsx_cell_style(cell: Cell, style: StyleDescriptor | None) -> None:
    if style is None:
        raise ValueError("XLSX cell style operation requires style")
    cell.font = _xlsx_font_with_style(cell, style)
    _apply_xlsx_fill_alignment_and_number_format(cell, style)


def _xlsx_font_with_style(cell: Cell, style: StyleDescriptor) -> Font:
    font = copy(cell.font)
    if style.bold is not None:
        font.bold = style.bold
    if style.italic is not None:
        font.italic = style.italic
    if style.underline is not None:
        font.underline = "single" if style.underline else None
    if style.font_family is not None:
        font.name = style.font_family
    if style.font_size_pt is not None:
        font.sz = float(style.font_size_pt)
    if style.font_color_rgb is not None:
        font.color = style.font_color_rgb
    return cast(Font, font)


def _apply_xlsx_fill_alignment_and_number_format(cell: Cell, style: StyleDescriptor) -> None:
    if style.fill_color_rgb is not None:
        cell.fill = PatternFill("solid", fgColor=style.fill_color_rgb)
    if style.alignment is not None:
        cell.alignment = Alignment(horizontal=style.alignment)
    if style.number_format is not None:
        cell.number_format = style.number_format


def _set_xlsx_metadata(workbook: openpyxl.Workbook, operation: DocumentPatchOperation) -> None:
    property_name = operation.target_path.rsplit("/", 1)[-1]
    if property_name not in {"creator", "title", "subject", "description", "keywords"}:
        raise ValueError(f"Unsupported XLSX metadata target: {operation.target_path}")
    setattr(workbook.properties, property_name, _string_value(operation.value))


def _cell_has_non_default_style(cell: object) -> bool:
    return bool(getattr(cell, "has_style", False))


def _xlsx_style_descriptor(cell: Cell, source_path: str) -> StyleDescriptor:
    font = cell.font
    fill = cell.fill
    fill_color = getattr(fill, "fgColor", None)
    fill_rgb = getattr(fill_color, "rgb", None)
    if isinstance(fill_rgb, str) and len(fill_rgb) == 8:
        fill_rgb = fill_rgb[-6:]
    return StyleDescriptor(
        style_id=f"xlsx-style-{source_path.strip('/').replace('/', '-')}",
        target_path=source_path,
        font_family=getattr(font, "name", None),
        font_size_pt=Decimal(str(getattr(font, "sz", 0))) if getattr(font, "sz", None) else None,
        bold=getattr(font, "bold", None),
        italic=getattr(font, "italic", None),
        fill_color_rgb=fill_rgb if isinstance(fill_rgb, str) and len(fill_rgb) == 6 else None,
        number_format=getattr(cell, "number_format", None),
    )


def _apply_pptx_operation(
    presentation: PptxPresentation,
    operation: DocumentPatchOperation,
) -> None:
    if operation.operation_type in {OperationType.replace_text, OperationType.set_field_value}:
        _set_pptx_text(presentation, operation.target_path, _string_value(operation.value))
        return
    if operation.operation_type is OperationType.set_table_cell:
        table, row_index, cell_index = _pptx_table_cell(presentation, operation.target_path)
        table.cell(row_index, cell_index).text = _string_value(operation.value)
        return
    if operation.operation_type is OperationType.set_document_metadata:
        _set_pptx_metadata(presentation, operation)
        return
    raise ValueError(f"Unsupported PPTX operation: {operation.operation_type.value}")


def _set_pptx_text(
    presentation: PptxPresentation,
    target_path: str,
    value: str,
) -> None:
    if target_path.endswith("/placeholders/title"):
        slide = _pptx_slide(presentation, target_path)
        if slide.shapes.title is None:
            raise ValueError(f"PPTX title placeholder not found: {target_path}")
        slide.shapes.title.text = value
        return
    match = _PPTX_SHAPE_TEXT_RE.match(target_path)
    if match is None:
        raise ValueError(f"Unsupported PPTX text target: {target_path}")
    slide = presentation.slides[int(match.group("slide")) - 1]
    shape = slide.shapes[int(match.group("shape")) - 1]
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"PPTX shape has no text frame: {target_path}")
    shape.text = value


def _pptx_slide(presentation: PptxPresentation, target_path: str) -> PptxSlide:
    match = re.match(r"^/slides/(?P<slide>\d+)/", target_path)
    if match is None:
        raise ValueError(f"Unsupported PPTX slide target: {target_path}")
    return cast(PptxSlide, presentation.slides[int(match.group("slide")) - 1])


def _pptx_table_cell(
    presentation: PptxPresentation,
    target_path: str,
) -> tuple[PptxTable, int, int]:
    match = _PPTX_TABLE_CELL_RE.match(target_path)
    if match is None:
        raise ValueError(f"Unsupported PPTX table target: {target_path}")
    slide = presentation.slides[int(match.group("slide")) - 1]
    table_ordinal = int(match.group("table"))
    table_shape_count = 0
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table_shape_count += 1
            if table_shape_count == table_ordinal:
                return (
                    cast(PptxTable, shape.table),
                    int(match.group("row")) - 1,
                    int(match.group("cell")) - 1,
                )
    raise ValueError(f"PPTX table not found: {target_path}")


def _set_pptx_metadata(presentation: PptxPresentation, operation: DocumentPatchOperation) -> None:
    property_name = operation.target_path.rsplit("/", 1)[-1]
    if property_name not in {"author", "category", "comments", "keywords", "subject", "title"}:
        raise ValueError(f"Unsupported PPTX core metadata target: {operation.target_path}")
    setattr(presentation.core_properties, property_name, _string_value(operation.value))


def _pptx_table_block(
    table: PptxTable,
    *,
    slide_index: int,
    table_index: int,
) -> TableBlock:
    cells: list[TableCell] = []
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            cells.append(
                TableCell(
                    row_index=row_index,
                    column_index=column_index,
                    text=cell.text,
                    source_path=(
                        f"/slides/{slide_index}/tables/{table_index}/rows/"
                        f"{row_index + 1}/cells/{column_index + 1}"
                    ),
                )
            )
    return TableBlock(
        block_id=f"pptx-slide-{slide_index:03d}-table-{table_index:03d}",
        source_path=f"/slides/{slide_index}/tables/{table_index}",
        cells=cells,
    )


def _pptx_core_metadata(presentation: PptxPresentation) -> dict[str, MetadataValue]:
    properties = presentation.core_properties
    metadata: dict[str, MetadataValue] = {
        "engine_id": "python-pptx",
        "format": "pptx",
        "slide_count": len(presentation.slides),
    }
    for property_name in (
        "author",
        "category",
        "comments",
        "created",
        "keywords",
        "last_modified_by",
        "modified",
        "revision",
        "subject",
        "title",
        "version",
    ):
        value = getattr(properties, property_name)
        if _metadata_value_is_present(value):
            metadata[f"core_{property_name}"] = value
    return metadata


def _slide_index_from_path(source_path: str) -> int:
    match = re.match(r"^/slides/(?P<slide>\d+)/", source_path)
    return int(match.group("slide")) if match is not None else 1


def _save_workbook_to_bytes(workbook: openpyxl.Workbook) -> bytes:
    from io import BytesIO

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _save_to_bytes(document: _OfficeSaveable) -> bytes:
    from io import BytesIO

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _svg_page(lines: list[str], *, title: str) -> bytes:
    escaped_title = html.escape(title)
    text_nodes = []
    for index, line in enumerate(lines[:40], start=1):
        text_nodes.append(
            f'<text x="48" y="{64 + index * 24}" font-family="Arial" '
            f'font-size="16">{html.escape(str(line))}</text>'
        )
    payload = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="960" height="1240" '
        'viewBox="0 0 960 1240">'
        '<rect width="960" height="1240" fill="#fff"/>'
        '<rect x="28" y="28" width="904" height="1184" fill="none" '
        'stroke="#c7c7c7" stroke-width="2"/>'
        f'<text x="48" y="54" font-family="Arial" font-size="20" '
        f'font-weight="700">{escaped_title}</text>'
        f"{''.join(text_nodes)}"
        "</svg>"
    )
    return payload.encode("utf-8")


def _string_value(value: ScalarValue) -> str:
    if value is None:
        return ""
    if isinstance(value, date | datetime):
        return value.isoformat()
    return str(value)


def _office_scalar(value: ScalarValue) -> str | int | float | bool | date | datetime | None:
    if isinstance(value, Decimal):
        return float(value)
    return value
