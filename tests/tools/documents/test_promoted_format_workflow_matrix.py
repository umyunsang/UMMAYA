# SPDX-License-Identifier: Apache-2.0
"""Promoted document format workflow matrix tests.

These tests exercise the single model-facing ``document`` primitive for each
write-promoted format.  A format is not considered practically promoted unless
the same flow can mutate, render, save to a user-visible path, and re-read the
saved derivative.
"""

from __future__ import annotations

import hashlib
from decimal import Decimal
from pathlib import Path

import pytest
from odfdo import Document as OdfDocument
from odfdo import DrawPage, Frame, Paragraph, Table

from tests.tools.documents.test_builtin_hwpx_engine import _write_rhwp_text_fixture
from tests.tools.documents.test_ooxml_adapters import _write_docx_template
from tests.tools.documents.test_pdf_adapter import _write_acroform_pdf
from ummaya.tools.documents.models import (
    DocumentExtraction,
    DocumentFormat,
    ToolResultStatus,
)
from ummaya.tools.documents.registry import DocumentToolRuntime
from ummaya.tools.documents.tool_defs import (
    DocumentFieldPatch,
    DocumentInspectRequest,
    DocumentLocator,
    DocumentPrimitiveRequest,
    DocumentStylePatch,
)


def test_hwpx_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = _write_rhwp_text_fixture(tmp_path / "weekly.hwpx")
    destination = tmp_path / "completed.hwpx"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.hwpx,
        source=source,
        destination=destination,
        target_path="/hwpx/text[1]",
        before_value="12 주차",
        after_value="13 주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "13 주차" in extraction_text


def test_hwpx_style_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = _write_rhwp_text_fixture(tmp_path / "styled.hwpx")
    destination = tmp_path / "styled-completed.hwpx"
    runtime = DocumentToolRuntime(
        session_id="promoted-matrix-hwpx-style",
        artifact_root=tmp_path / "store-hwpx-style",
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="promoted-matrix-hwpx-style",
            document=DocumentLocator(
                path=str(source),
                expected_format=DocumentFormat.hwpx,
            ),
            operation="style",
            instruction=(
                "Apply HWPX public-form style proof and save the reviewed "
                f"derivative to {destination}."
            ),
            styles=(
                DocumentStylePatch(
                    target_path="/hwpx/text[1]",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("14"),
                    bold=True,
                    font_color_rgb="1F4E79",
                    fill_color_rgb="FFF2CC",
                    alignment="right",
                ),
            ),
            destination_display_name=destination.name,
            destination_path=str(destination),
        )
    )

    assert result.status is ToolResultStatus.ok
    assert result.diff is not None
    assert [change.change_type for change in result.diff.changes] == ["style"]
    assert result.saved_exports
    saved = result.saved_exports[0]
    assert saved.local_path == destination.resolve()
    assert destination.is_file()
    assert saved.sha256 == hashlib.sha256(destination.read_bytes()).hexdigest()
    assert saved.byte_size == destination.stat().st_size
    assert _workflow_status(result, "save") == "completed"
    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")

    reread = runtime.inspect(
        DocumentInspectRequest(
            correlation_id="promoted-matrix-hwpx-style-reread",
            document=DocumentLocator(
                path=str(destination),
                expected_format=DocumentFormat.hwpx,
            ),
        )
    )

    assert reread.status is ToolResultStatus.ok
    assert reread.extraction is not None
    assert any(
        style.font_family == "Malgun Gothic"
        and style.font_size_pt == Decimal("14")
        and style.bold is True
        and style.font_color_rgb == "1F4E79"
        and style.fill_color_rgb == "FFF2CC"
        for style in reread.extraction.style_map
    )
    assert any(style.alignment == "right" for style in reread.extraction.style_map)


def test_owpml_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = _write_rhwp_text_fixture(tmp_path / "weekly.owpml")
    destination = tmp_path / "completed.owpml"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.owpml,
        source=source,
        destination=destination,
        target_path="/hwpx/text[1]",
        before_value="12 주차",
        after_value="13 주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "13 주차" in extraction_text


def test_docx_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = _write_docx_template(tmp_path / "civil-form.docx")
    destination = tmp_path / "completed.docx"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.docx,
        source=source,
        destination=destination,
        target_path="engine://python-docx/civil-form.docx/table/1/r1c2",
        before_value="2026-06-01",
        after_value="2026-06-03",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "2026-06-03" in extraction_text


def test_xlsx_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    source = tmp_path / "ledger.xlsx"
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "제출서류"
    worksheet["A1"] = "주차"
    worksheet["B1"] = "13주차"
    workbook.save(source)
    destination = tmp_path / "completed.xlsx"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.xlsx,
        source=source,
        destination=destination,
        target_path="/sheets/제출서류/cells/B1",
        before_value="13주차",
        after_value="14주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차" in extraction_text


def test_pptx_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    source = tmp_path / "briefing.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "항목"
    table_shape.table.cell(0, 1).text = "13주차"
    presentation.save(source)
    destination = tmp_path / "completed.pptx"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.pptx,
        source=source,
        destination=destination,
        target_path="/slides/1/tables/1/rows/1/cells/2",
        before_value="13주차",
        after_value="14주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차" in extraction_text


def test_pdf_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = _write_acroform_pdf(tmp_path / "application.pdf")
    destination = tmp_path / "completed.pdf"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.pdf,
        source=source,
        destination=destination,
        target_path="/acroform/fields/applicant_name",
        before_value="",
        after_value="Hong Gil Dong",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/png")
    assert "Hong Gil Dong" in extraction_text


def test_odt_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "notice.odt"
    document = OdfDocument("text")
    document.body.clear()
    document.body.append(Paragraph("13주차"))
    document.save(source)
    destination = tmp_path / "completed.odt"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat("odt"),
        source=source,
        destination=destination,
        target_path="/odf/text/p[1]",
        before_value="13주차",
        after_value="14주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차" in extraction_text


def test_ods_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "ledger.ods"
    document = OdfDocument("spreadsheet")
    document.body.clear()
    table = Table("제출서류")
    table.set_value("A1", "주차")
    table.set_value("B1", "13주차")
    document.body.append(table)
    document.save(source)
    destination = tmp_path / "completed.ods"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat("ods"),
        source=source,
        destination=destination,
        target_path="/odf/sheets/제출서류/cells/B1",
        before_value="13주차",
        after_value="14주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차" in extraction_text


def test_odp_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "briefing.odp"
    document = OdfDocument("presentation")
    document.body.clear()
    page = DrawPage("page1", name="Page 1")
    page.append(
        Frame.text_frame(
            "13주차",
            size=("7cm", "5cm"),
            position=("1cm", "1cm"),
            style="Standard",
            text_style="Standard",
        )
    )
    document.body.append(page)
    document.save(source)
    destination = tmp_path / "completed.odp"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat("odp"),
        source=source,
        destination=destination,
        target_path="/odf/slides/1/frames/1",
        before_value="13주차",
        after_value="14주차",
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차" in extraction_text


@pytest.mark.parametrize(
    ("suffix", "source_text", "before_text", "after_text"),
    (
        ("txt", "13주차 활동일지\n", "13주차 활동일지", "14주차 활동일지\n"),
        ("md", "# 13주차 활동일지\n", "# 13주차 활동일지", "# 14주차 활동일지\n"),
        (
            "html",
            "<html><body><p>13주차 활동일지</p></body></html>\n",
            "13주차 활동일지",
            "14주차 활동일지\n",
        ),
        (
            "htm",
            "<html><body><p>13주차 활동일지</p></body></html>\n",
            "13주차 활동일지",
            "14주차 활동일지\n",
        ),
        ("rtf", r"{\rtf1\ansi 13주차 활동일지}" + "\n", "13주차 활동일지", "14주차 활동일지\n"),
    ),
)
def test_text_web_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
    suffix: str,
    source_text: str,
    before_text: str,
    after_text: str,
) -> None:
    source = tmp_path / f"notice.{suffix}"
    source.write_text(source_text, encoding="utf-8")
    destination = tmp_path / f"completed.{suffix}"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat(suffix),
        source=source,
        destination=destination,
        target_path="/text/body",
        before_value=before_text,
        after_value=after_text,
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14주차 활동일지" in extraction_text


@pytest.mark.parametrize(
    ("suffix", "source_text", "before_text", "after_text"),
    (
        ("csv", "week,value\n13,old\n", "week,value\n13,old", "week,value\n14,new\n"),
        ("tsv", "week\tvalue\n13\told\n", "week\tvalue\n13\told", "week\tvalue\n14\tnew\n"),
        (
            "json",
            '{"week": 13, "value": "old"}\n',
            '{"value": "old", "week": 13}',
            '{"week": 14, "value": "new"}\n',
        ),
        ("jsonl", '{"week":13}\n', '{"week": 13}', '{"week":14}\n'),
        ("yaml", "week: 13\nvalue: old\n", "value: old\nweek: 13", "week: 14\nvalue: new\n"),
        ("yml", "week: 13\nvalue: old\n", "value: old\nweek: 13", "week: 14\nvalue: new\n"),
        ("xml", "<root><week>13</week></root>\n", "13", "<root><week>14</week></root>\n"),
        ("rdf", "<rdf><week>13</week></rdf>\n", "13", "<rdf><week>14</week></rdf>\n"),
        ("gpx", "<gpx><name>13</name></gpx>\n", "13", "<gpx><name>14</name></gpx>\n"),
        ("kml", "<kml><name>13</name></kml>\n", "13", "<kml><name>14</name></kml>\n"),
        ("hml", "<hml><week>13</week></hml>\n", "13", "<hml><week>14</week></hml>\n"),
        (
            "ttl",
            '@prefix ex: <https://example.test/> .\nex:w ex:v "13" .\n',
            '@prefix ex: <https://example.test/> .\nex:w ex:v "13" .',
            '@prefix ex: <https://example.test/> .\nex:w ex:v "14" .\n',
        ),
        (
            "lod",
            '<https://example.test/w> <https://example.test/v> "13" .\n',
            '<https://example.test/w> <https://example.test/v> "13" .',
            '<https://example.test/w> <https://example.test/v> "14" .\n',
        ),
        ("fasta", ">week\nACTG13\n", ">week\nACTG13", ">week\nACTG14\n"),
        (
            "sgml",
            "<doc><week>13</week></doc>\n",
            "<doc><week>13</week></doc>",
            "<doc><week>14</week></doc>\n",
        ),
        (
            "dtd",
            "<!ELEMENT week (#PCDATA)>\n",
            "<!ELEMENT week (#PCDATA)>",
            "<!ELEMENT week (#PCDATA)>\n",
        ),
        ("etc", "week=13\n", "week=13", "week=14\n"),
    ),
)
def test_data_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
    suffix: str,
    source_text: str,
    before_text: str,
    after_text: str,
) -> None:
    source = tmp_path / f"dataset.{suffix}"
    source.write_text(source_text, encoding="utf-8")
    destination = tmp_path / f"completed.{suffix}"

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat(suffix),
        source=source,
        destination=destination,
        target_path="/data/body",
        before_value=before_text,
        after_value=after_text,
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert "14" in extraction_text or suffix == "dtd"


def test_python_source_document_primitive_save_renders_rereads_and_diffs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "script.py"
    source.write_text('def main():\n    return "13주차"\n', encoding="utf-8")
    destination = tmp_path / "completed.py"
    replacement = 'def main():\n    return "14주차"\n'

    result, extraction_text = _run_save_flow(
        tmp_path,
        document_format=DocumentFormat.python,
        source=source,
        destination=destination,
        target_path="/code/body",
        before_value='def main():\n    return "13주차"',
        after_value=replacement,
    )

    _assert_render_artifacts_exist(result, expected_mime_type="image/svg+xml")
    assert 'return "14주차"' in extraction_text


def _run_save_flow(
    tmp_path: Path,
    *,
    document_format: DocumentFormat,
    source: Path,
    destination: Path,
    target_path: str,
    before_value: str,
    after_value: str,
) -> tuple[object, str]:
    runtime = DocumentToolRuntime(
        session_id=f"promoted-matrix-{document_format.value}",
        artifact_root=tmp_path / f"store-{document_format.value}",
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id=f"promoted-matrix-{document_format.value}",
            document=DocumentLocator(
                path=str(source),
                expected_format=document_format,
            ),
            operation="save",
            instruction=(
                "Complete this public document fixture and save the reviewed "
                f"derivative to {destination}."
            ),
            patches=(DocumentFieldPatch(target_path=target_path, value=after_value),),
            destination_display_name=destination.name,
            destination_path=str(destination),
        )
    )

    assert result.status is ToolResultStatus.ok
    assert result.diff is not None
    assert [(change.before_value, change.after_value) for change in result.diff.changes] == [
        (before_value, after_value)
    ]
    assert result.saved_exports
    saved = result.saved_exports[0]
    assert saved.local_path == destination.resolve()
    assert destination.is_file()
    assert saved.sha256 == hashlib.sha256(destination.read_bytes()).hexdigest()
    assert saved.byte_size == destination.stat().st_size
    assert _workflow_status(result, "save") == "completed"

    reread = runtime.inspect(
        DocumentInspectRequest(
            correlation_id=f"promoted-matrix-{document_format.value}-reread",
            document=DocumentLocator(
                path=str(destination),
                expected_format=document_format,
            ),
        )
    )
    assert reread.status is ToolResultStatus.ok
    assert reread.extraction is not None
    return result, _flatten_extraction_text(reread.extraction)


def _assert_render_artifacts_exist(result: object, *, expected_mime_type: str) -> None:
    render_artifacts = result.render_artifacts
    assert render_artifacts
    first_render = render_artifacts[0]
    render_path = Path(first_render.render_path)
    assert first_render.render_mime_type == expected_mime_type
    assert render_path.is_file()
    payload = render_path.read_bytes()
    if expected_mime_type == "image/png":
        assert payload.startswith(b"\x89PNG\r\n\x1a\n")
    else:
        assert payload.lstrip().startswith(b"<svg")


def _workflow_status(result: object, step_id: str) -> str | None:
    for step in result.workflow_steps:
        if step.step_id == step_id:
            return step.status
    return None


def _flatten_extraction_text(extraction: DocumentExtraction) -> str:
    chunks: list[str] = []
    chunks.extend(paragraph.text for paragraph in extraction.paragraphs)
    for table in extraction.tables:
        chunks.extend(cell.text for cell in table.cells)
    chunks.extend(str(field.current_value) for field in extraction.fields)
    chunks.extend(str(value) for value in extraction.metadata.values())
    return "\n".join(chunks)
