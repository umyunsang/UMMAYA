# SPDX-License-Identifier: Apache-2.0
"""OOXML adapter boundary tests for DOCX, XLSX, and PPTX."""

from __future__ import annotations

import base64
import zipfile
from decimal import Decimal
from pathlib import Path

import docx
import pytest

from ummaya.tools.documents.adapter_registry import (
    UnsupportedDocumentAdapterError,
    build_default_document_adapter_registry,
    build_document_adapter_registry_from_engine_registry,
)
from ummaya.tools.documents.artifact_store import DocumentArtifactStore
from ummaya.tools.documents.engines import DocumentEngineRegistry
from ummaya.tools.documents.formats.ooxml import (
    DocxDocumentAdapter,
    OpenPyxlDocumentEngine,
    PptxDocumentAdapter,
    PythonDocxDocumentEngine,
    PythonPptxDocumentEngine,
    XlsxDocumentAdapter,
)
from ummaya.tools.documents.models import (
    DocumentArtifact,
    DocumentExtraction,
    DocumentFormat,
    DocumentIR,
    DocumentPatch,
    DocumentPatchOperation,
    FormSlot,
    KnownDocumentFormat,
    OperationType,
    StyleDescriptor,
    TableCell,
    ToolResultStatus,
)
from ummaya.tools.documents.patch import apply_document_patch, copy_for_edit
from ummaya.tools.documents.registry import _DOCX_TABLE_FILL_TARGET_RE, DocumentToolRuntime
from ummaya.tools.documents.tool_defs import (
    DocumentFieldPatch,
    DocumentLocator,
    DocumentPrimitiveRequest,
    DocumentStylePatch,
)

_WEB_FORM_SOURCE_DIR = Path(".evidence/public-document-design-web-forms-20260604/sources")
_SEOUL_CULTURE_DOCX = _WEB_FORM_SOURCE_DIR / "seoul-culture-application-plan.docx"


def test_default_registry_uses_separate_ooxml_adapters() -> None:
    registry = build_default_document_adapter_registry()

    docx_adapter = registry.require_promoted(DocumentFormat.docx)
    xlsx_adapter = registry.require_promoted(DocumentFormat.xlsx)
    pptx_adapter = registry.require_promoted(DocumentFormat.pptx)

    assert isinstance(docx_adapter, DocxDocumentAdapter)
    assert isinstance(xlsx_adapter, XlsxDocumentAdapter)
    assert isinstance(pptx_adapter, PptxDocumentAdapter)
    assert docx_adapter.adapter_id == "python-docx-adapter"
    assert xlsx_adapter.adapter_id == "openpyxl-adapter"
    assert pptx_adapter.adapter_id == "python-pptx-adapter"


def test_ooxml_registry_can_keep_unpromoted_known_adapters() -> None:
    registry = build_document_adapter_registry_from_engine_registry(DocumentEngineRegistry())

    assert isinstance(registry.require_known(KnownDocumentFormat.xlsx), XlsxDocumentAdapter)
    assert isinstance(registry.require_known(KnownDocumentFormat.pptx), PptxDocumentAdapter)

    with pytest.raises(UnsupportedDocumentAdapterError):
        registry.require_promoted(DocumentFormat.xlsx)
    with pytest.raises(UnsupportedDocumentAdapterError):
        registry.require_promoted(DocumentFormat.pptx)


def test_docx_engine_writes_paragraph_run_table_metadata_and_preserves_style(
    tmp_path: Path,
) -> None:
    source = _write_docx_template(tmp_path / "civil-form.docx")
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.docx,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonDocxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-docx",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="paragraph-title",
                operation_type=OperationType.replace_text,
                target_path="/word/paragraphs/1",
                value="공공 AX 신청서",
            ),
            DocumentPatchOperation(
                operation_id="run-applicant",
                operation_type=OperationType.replace_text,
                target_path="/word/paragraphs/2/runs/2",
                value="김철수",
            ),
            DocumentPatchOperation(
                operation_id="table-cell",
                operation_type=OperationType.set_table_cell,
                target_path="/word/tables/1/rows/1/cells/2",
                value="2026-06-03",
            ),
            DocumentPatchOperation(
                operation_id="metadata-title",
                operation_type=OperationType.set_document_metadata,
                target_path="/core/title",
                value="수정된 공공 AX 신청서",
            ),
            DocumentPatchOperation(
                operation_id="run-style",
                operation_type=OperationType.set_run_style,
                target_path="/word/paragraphs/2/runs/2",
                style=StyleDescriptor(
                    style_id="run-applicant-style",
                    target_path="/word/paragraphs/2/runs/2",
                    bold=True,
                    font_size_pt=Decimal("12"),
                ),
            ),
        ],
        dry_run=False,
        expected_format=DocumentFormat.docx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-docx",
        destination_name="derivative.docx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    updated = docx.Document(str(result.derivative_artifact.source_path))
    assert updated.paragraphs[0].text == "공공 AX 신청서"
    assert updated.paragraphs[0].style is not None
    assert updated.paragraphs[0].style.name == "Title"
    assert updated.paragraphs[1].runs[1].text == "김철수"
    assert updated.paragraphs[1].runs[1].bold is True
    assert updated.tables[0].cell(0, 1).text == "2026-06-03"
    assert updated.core_properties.title == "수정된 공공 AX 신청서"


def test_docx_run_style_writes_east_asia_font_attrs(tmp_path: Path) -> None:
    source = _write_docx_template(tmp_path / "civil-form.docx")
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.docx,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonDocxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-docx-eastasia-run",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="run-style-eastasia",
                operation_type=OperationType.set_run_style,
                target_path="/word/paragraphs/2/runs/2",
                style=StyleDescriptor(
                    style_id="run-style-eastasia",
                    target_path="/word/paragraphs/2/runs/2",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("12"),
                    bold=True,
                    font_color_rgb="1F4E79",
                ),
            )
        ],
        dry_run=False,
        expected_format=DocumentFormat.docx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-docx-eastasia-run",
        destination_name="derivative-eastasia-run.docx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    document_xml = _docx_xml(result.derivative_artifact.source_path)
    assert 'w:ascii="Malgun Gothic"' in document_xml
    assert 'w:hAnsi="Malgun Gothic"' in document_xml
    assert 'w:eastAsia="Malgun Gothic"' in document_xml
    assert 'w:cs="Malgun Gothic"' in document_xml
    styles_xml = _docx_xml(result.derivative_artifact.source_path, "word/styles.xml")
    assert 'w:eastAsia="Malgun Gothic"' not in styles_xml


def test_docx_paragraph_style_writes_all_font_attrs_on_affected_runs(
    tmp_path: Path,
) -> None:
    source = _write_docx_template(tmp_path / "civil-form-paragraph.docx")
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.docx,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonDocxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-docx-paragraph-font-parity",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="paragraph-style-font-parity",
                operation_type=OperationType.set_paragraph_style,
                target_path="/word/paragraphs/2",
                style=StyleDescriptor(
                    style_id="paragraph-style-font-parity",
                    target_path="/word/paragraphs/2",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("12"),
                    bold=True,
                    font_color_rgb="1F4E79",
                ),
            )
        ],
        dry_run=False,
        expected_format=DocumentFormat.docx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-docx-paragraph-font-parity",
        destination_name="derivative-paragraph-font-parity.docx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    document_xml = _docx_xml(result.derivative_artifact.source_path)
    assert document_xml.count('w:ascii="Malgun Gothic"') == 2
    assert document_xml.count('w:hAnsi="Malgun Gothic"') == 2
    assert document_xml.count('w:eastAsia="Malgun Gothic"') == 2
    assert document_xml.count('w:cs="Malgun Gothic"') == 2
    styles_xml = _docx_xml(result.derivative_artifact.source_path, "word/styles.xml")
    assert 'w:eastAsia="Malgun Gothic"' not in styles_xml


def test_docx_cell_style_writes_east_asia_font_attrs_on_public_form(tmp_path: Path) -> None:
    source = _write_docx_template(tmp_path / "civil-form-public.docx")
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.docx,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonDocxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-docx-eastasia-cell",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="cell-style-eastasia",
                operation_type=OperationType.set_cell_style,
                target_path="/word/tables/1/rows/1/cells/2",
                style=StyleDescriptor(
                    style_id="cell-style-eastasia",
                    target_path="/word/tables/1/rows/1/cells/2",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("11"),
                    fill_color_rgb="FFF2CC",
                    alignment="center",
                ),
            )
        ],
        dry_run=False,
        expected_format=DocumentFormat.docx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-docx-eastasia-cell",
        destination_name="derivative-eastasia-cell.docx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    document_xml = _docx_xml(result.derivative_artifact.source_path)
    assert 'w:fill="FFF2CC"' in document_xml
    assert 'w:jc w:val="center"' in document_xml
    assert 'w:ascii="Malgun Gothic"' in document_xml
    assert 'w:hAnsi="Malgun Gothic"' in document_xml
    assert 'w:eastAsia="Malgun Gothic"' in document_xml
    assert 'w:cs="Malgun Gothic"' in document_xml
    styles_xml = _docx_xml(result.derivative_artifact.source_path, "word/styles.xml")
    assert 'w:eastAsia="Malgun Gothic"' not in styles_xml


def test_docx_style_without_font_family_does_not_invent_eastasia_attr(tmp_path: Path) -> None:
    source = _write_docx_template(tmp_path / "civil-form-no-font-family.docx")
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.docx,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonDocxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-docx-no-font-family",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="run-style-without-font-family",
                operation_type=OperationType.set_run_style,
                target_path="/word/paragraphs/2/runs/2",
                style=StyleDescriptor(
                    style_id="run-style-without-font-family",
                    target_path="/word/paragraphs/2/runs/2",
                    bold=True,
                    font_size_pt=Decimal("12"),
                ),
            )
        ],
        dry_run=False,
        expected_format=DocumentFormat.docx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-docx-no-font-family",
        destination_name="derivative-no-font-family.docx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    document_xml = _docx_xml(result.derivative_artifact.source_path)
    assert "w:eastAsia=" not in document_xml
    assert "w:cs=" not in document_xml


def test_docx_public_form_extracts_adjacent_label_blank_cell_slots() -> None:
    assert _SEOUL_CULTURE_DOCX.exists()
    engine = PythonDocxDocumentEngine()
    extraction = engine.inspect(_SEOUL_CULTURE_DOCX, artifact_id="seoul-culture-application-plan")
    ir = DocumentIR.from_extraction(
        artifact_id="seoul-culture-application-plan",
        document_format=DocumentFormat.docx,
        extraction=extraction,
        engine_id=engine.engine_id,
    )

    team_slot = _slot_by_label(ir, "팀명")
    brand_slot = _slot_by_label(ir, "기업 및 브랜드명")
    receipt_slot = _slot_by_label(ir, "접수번호")

    assert team_slot.current_value == ""
    assert team_slot.source_anchor.format_path == (
        "engine://python-docx/seoul-culture-application-plan.docx/table/1/r3c2"
    )
    assert _DOCX_TABLE_FILL_TARGET_RE.search(team_slot.source_anchor.format_path) is not None
    assert _table_cell_by_path(extraction, team_slot.source_anchor.format_path).field_path == (
        team_slot.source_anchor.format_path
    )
    assert brand_slot.current_value == ""
    assert _DOCX_TABLE_FILL_TARGET_RE.search(brand_slot.source_anchor.format_path) is not None
    assert _table_cell_by_path(extraction, brand_slot.source_anchor.format_path).field_path == (
        brand_slot.source_anchor.format_path
    )
    assert receipt_slot.current_value == ""
    assert receipt_slot.source_anchor.format_path == (
        "engine://python-docx/seoul-culture-application-plan.docx/table/1/r1c7"
    )
    assert _DOCX_TABLE_FILL_TARGET_RE.search(receipt_slot.source_anchor.format_path) is not None
    assert _table_cell_by_path(extraction, receipt_slot.source_anchor.format_path).field_path == (
        receipt_slot.source_anchor.format_path
    )


def test_openpyxl_engine_writes_cells_styles_and_preserves_workbook_contract(
    tmp_path: Path,
) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.styles import Font, PatternFill  # type: ignore[import-untyped]

    source = tmp_path / "ledger.xlsx"
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "제출서류"
    worksheet["A1"] = "신청인"
    worksheet["B1"] = "홍길동"
    worksheet["C1"] = "=CONCAT(A1,B1)"
    worksheet["B2"] = "병합"
    worksheet.merge_cells("B2:C2")
    worksheet["B2"].font = Font(bold=True)
    worksheet["B2"].fill = PatternFill("solid", fgColor="FFFF00")
    worksheet.print_area = "A1:C10"
    workbook.save(source)
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.xlsx,
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )
    registry = DocumentEngineRegistry()
    registry.register(OpenPyxlDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-xlsx-real",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="set-b1",
                operation_type=OperationType.set_table_cell,
                target_path="/sheets/제출서류/cells/B1",
                value="김철수",
            ),
            DocumentPatchOperation(
                operation_id="style-b2",
                operation_type=OperationType.set_cell_style,
                target_path="/sheets/제출서류/cells/B2",
                style=StyleDescriptor(
                    style_id="style-b2",
                    target_path="/sheets/제출서류/cells/B2",
                    fill_color_rgb="00FF00",
                    number_format="@",
                    bold=True,
                ),
            ),
        ],
        dry_run=False,
        expected_format=DocumentFormat.xlsx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-xlsx",
        destination_name="derivative.xlsx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    reloaded = openpyxl.load_workbook(result.derivative_artifact.source_path, data_only=False)
    sheet = reloaded["제출서류"]
    assert sheet["B1"].value == "김철수"
    assert sheet["C1"].value == "=CONCAT(A1,B1)"
    assert str(sheet.print_area) == "'제출서류'!$A$1:$C$10"
    assert "B2:C2" in {str(region) for region in sheet.merged_cells.ranges}
    assert sheet["B2"].font.bold is True
    assert sheet["B2"].number_format == "@"
    assert sheet["B2"].fill.fgColor.rgb == "0000FF00"


def test_python_pptx_engine_writes_placeholders_tables_metadata_and_blocks_media(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    source = tmp_path / "briefing.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "기존 제목"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(0.5))
    textbox.text = "본문"
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "항목"
    table_shape.table.cell(0, 1).text = "값"
    image_path = tmp_path / "pixel.png"
    image_path.write_bytes(base64.b64decode(_ONE_PIXEL_PNG_BASE64))
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3), Inches(1), Inches(1))
    presentation.core_properties.title = "공공 AX 발표"
    presentation.save(source)
    store, working = _working_artifact(
        tmp_path,
        source,
        document_format=DocumentFormat.pptx,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    registry = DocumentEngineRegistry()
    registry.register(PythonPptxDocumentEngine())
    patch = DocumentPatch(
        patch_id="patch-pptx",
        target_artifact_id=working.artifact_id,
        operations=[
            DocumentPatchOperation(
                operation_id="title",
                operation_type=OperationType.replace_text,
                target_path="/slides/1/placeholders/title",
                value="수정된 제목",
            ),
            DocumentPatchOperation(
                operation_id="table",
                operation_type=OperationType.set_table_cell,
                target_path="/slides/1/tables/1/rows/1/cells/2",
                value="수정값",
            ),
            DocumentPatchOperation(
                operation_id="metadata",
                operation_type=OperationType.set_document_metadata,
                target_path="/core/title",
                value="수정된 공공 AX 발표",
            ),
        ],
        dry_run=False,
        expected_format=DocumentFormat.pptx,
        destination_policy="working_copy",
    )

    result = apply_document_patch(
        store,
        working,
        patch,
        engine_registry=registry,
        artifact_id="derivative-pptx",
        destination_name="derivative.pptx",
    )

    assert result.status is ToolResultStatus.ok
    assert result.derivative_artifact is not None
    updated = pptx.Presentation(str(result.derivative_artifact.source_path))
    assert updated.slides[0].shapes.title.text == "수정된 제목"
    assert updated.slides[0].shapes[2].table.cell(0, 1).text == "수정값"
    assert updated.core_properties.title == "수정된 공공 AX 발표"
    assert any(shape.shape_type == 13 for shape in updated.slides[0].shapes)

    blocked_patch = patch.model_copy(
        update={
            "patch_id": "patch-pptx-media",
            "operations": [
                DocumentPatchOperation(
                    operation_id="media",
                    operation_type=OperationType.replace_text,
                    target_path="/slides/1/media/1",
                    value="unsupported",
                )
            ],
        }
    )
    blocked = apply_document_patch(
        store,
        working,
        blocked_patch,
        engine_registry=registry,
        artifact_id="blocked-pptx",
        destination_name="blocked.pptx",
    )
    assert blocked.status is ToolResultStatus.blocked
    assert blocked.derivative_artifact is None


def test_document_primitive_maps_docx_table_patch_to_table_cell_operation(
    tmp_path: Path,
) -> None:
    source = _write_docx_template(tmp_path / "civil-form.docx")
    runtime = DocumentToolRuntime(
        session_id="docx-table-primitive", artifact_root=tmp_path / "store"
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="docx-table-primitive",
            document=DocumentLocator(path=str(source), expected_format=DocumentFormat.docx),
            operation="fill",
            instruction="테스트 값으로 작성해.",
            patches=(
                DocumentFieldPatch(
                    target_path="engine://python-docx/civil-form.docx/table/1/r1c2",
                    value="2026-06-03",
                ),
            ),
        )
    )

    assert result.status is ToolResultStatus.ok
    assert result.diff is not None
    assert [
        (change.target_path, change.before_value, change.after_value)
        for change in result.diff.changes
    ] == [
        (
            "engine://python-docx/civil-form.docx/table/1/r1c2",
            "2026-06-01",
            "2026-06-03",
        )
    ]


def test_document_primitive_maps_xlsx_cell_patch_to_table_cell_operation(
    tmp_path: Path,
) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    source = tmp_path / "ledger.xlsx"
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "제출서류"
    worksheet["A1"] = "주차"
    worksheet["B1"] = "13주차"
    workbook.save(source)
    runtime = DocumentToolRuntime(
        session_id="xlsx-cell-primitive", artifact_root=tmp_path / "store"
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="xlsx-cell-primitive",
            document=DocumentLocator(path=str(source), expected_format=DocumentFormat.xlsx),
            operation="fill",
            instruction="테스트 값으로 작성해.",
            patches=(DocumentFieldPatch(target_path="/sheets/제출서류/cells/B1", value="14주차"),),
        )
    )

    assert result.status is ToolResultStatus.ok
    assert result.diff is not None
    assert [
        (change.target_path, change.before_value, change.after_value)
        for change in result.diff.changes
    ] == [("/sheets/제출서류/cells/B1", "13주차", "14주차")]


def test_document_primitive_maps_docx_run_style_to_run_style_operation(
    tmp_path: Path,
) -> None:
    source = _write_docx_template(tmp_path / "civil-form.docx")
    runtime = DocumentToolRuntime(
        session_id="docx-run-style-primitive",
        artifact_root=tmp_path / "store",
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="docx-run-style-primitive",
            document=DocumentLocator(path=str(source), expected_format=DocumentFormat.docx),
            operation="style",
            instruction="신청인 값의 글꼴, 크기, 굵기를 공문서 테스트 서식으로 바꿔줘.",
            styles=(
                DocumentStylePatch(
                    target_path="/word/paragraphs/2/runs/2",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("14"),
                    bold=True,
                    font_color_rgb="005BAC",
                ),
            ),
        )
    )

    assert result.status is ToolResultStatus.ok
    derivative_id = next(ref for ref in result.artifact_refs if ref.startswith("derivative-"))
    updated = docx.Document(str(runtime._artifacts[derivative_id].source_path))
    applicant_run = updated.paragraphs[1].runs[1]
    assert applicant_run.font.name == "Malgun Gothic"
    assert applicant_run.font.size is not None
    assert applicant_run.font.size.pt == 14
    assert applicant_run.bold is True
    assert str(applicant_run.font.color.rgb) == "005BAC"


def test_document_primitive_maps_xlsx_cell_style_to_cell_style_operation(
    tmp_path: Path,
) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    source = tmp_path / "summary.xlsx"
    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "총괄표"
    worksheet["B5"] = "기관명"
    workbook.save(source)
    runtime = DocumentToolRuntime(
        session_id="xlsx-cell-style-primitive",
        artifact_root=tmp_path / "store",
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="xlsx-cell-style-primitive",
            document=DocumentLocator(path=str(source), expected_format=DocumentFormat.xlsx),
            operation="style",
            instruction="총괄표 B5 셀의 글꼴, 배경색, 정렬을 공문서 테스트 서식으로 바꿔줘.",
            styles=(
                DocumentStylePatch(
                    target_path="/sheets/총괄표/cells/B5",
                    font_family="Malgun Gothic",
                    font_size_pt=Decimal("11"),
                    bold=True,
                    fill_color_rgb="FFF2CC",
                    alignment="center",
                ),
            ),
        )
    )

    assert result.status is ToolResultStatus.ok
    derivative_id = next(ref for ref in result.artifact_refs if ref.startswith("derivative-"))
    reloaded = openpyxl.load_workbook(runtime._artifacts[derivative_id].source_path)
    cell = reloaded["총괄표"]["B5"]
    assert cell.font.name == "Malgun Gothic"
    assert cell.font.sz == 11
    assert cell.font.bold is True
    assert cell.fill.fgColor.rgb == "00FFF2CC"
    assert cell.alignment.horizontal == "center"


def test_document_primitive_maps_pptx_table_patch_to_table_cell_operation(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    source = tmp_path / "briefing.pptx"
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "항목"
    table_shape.table.cell(0, 1).text = "13주차"
    presentation.save(source)
    runtime = DocumentToolRuntime(
        session_id="pptx-table-primitive", artifact_root=tmp_path / "store"
    )

    result = runtime.document(
        DocumentPrimitiveRequest(
            correlation_id="pptx-table-primitive",
            document=DocumentLocator(path=str(source), expected_format=DocumentFormat.pptx),
            operation="fill",
            instruction="테스트 값으로 작성해.",
            patches=(
                DocumentFieldPatch(
                    target_path="/slides/1/tables/1/rows/1/cells/2",
                    value="14주차",
                ),
            ),
        )
    )

    assert result.status is ToolResultStatus.ok
    assert result.diff is not None
    assert [
        (change.target_path, change.before_value, change.after_value)
        for change in result.diff.changes
    ] == [("/slides/1/tables/1/rows/1/cells/2", "13주차", "14주차")]


def _working_artifact(
    tmp_path: Path,
    source: Path,
    *,
    document_format: DocumentFormat,
    mime_type: str,
) -> tuple[DocumentArtifactStore, DocumentArtifact]:
    store = DocumentArtifactStore(root=tmp_path / f"store-{document_format.value}", session_id="s")
    source_artifact = store.store_source(
        source,
        artifact_id=f"source-{document_format.value}",
        document_format=document_format,
        mime_type=mime_type,
    )
    return store, copy_for_edit(
        store,
        source_artifact,
        artifact_id=f"working-{document_format.value}",
        destination_name=f"working.{document_format.value}",
    )


def _slot_by_label(document_ir: DocumentIR, label: str) -> FormSlot:
    for slot in document_ir.form_slots:
        if " ".join(slot.label.split()) == label:
            return slot
    raise KeyError(label)


def _table_cell_by_path(extraction: DocumentExtraction, source_path: str) -> TableCell:
    for table in extraction.tables:
        for cell in table.cells:
            if cell.source_path == source_path:
                return cell
    raise KeyError(source_path)


def _write_docx_template(path: Path) -> Path:
    document = docx.Document()
    title = document.add_paragraph("기존 신청서")
    title.style = "Title"
    paragraph = document.add_paragraph()
    paragraph.add_run("신청인: ")
    applicant_run = paragraph.add_run("홍길동")
    applicant_run.bold = False
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "작성일"
    table.cell(0, 1).text = "2026-06-01"
    document.core_properties.title = "기존 공공 AX 신청서"
    document.save(str(path))
    return path


def _docx_xml(path: Path, part_name: str = "word/document.xml") -> str:
    with zipfile.ZipFile(path) as archive:
        return archive.read(part_name).decode("utf-8")


_ONE_PIXEL_PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMB/6X6x9QAAAAASUVORK5CYII="
)
